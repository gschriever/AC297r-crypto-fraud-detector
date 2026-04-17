from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation()

# Slide 1: Unsupervised Learning Rationale
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Unsupervised Learning Rationale"
tf = slide.shapes.placeholders[1].text_frame
p = tf.add_paragraph()
p.text = "Class Imbalance: The vast majority of tokens are non-fraudulent, making it difficult to train a balanced supervised model."
p.level = 0
p = tf.add_paragraph()
p.text = "Labeling Bottleneck: It is practically impossible to manually label 5,000 mid-cap tokens accurately within our timeline."
p.level = 0
p = tf.add_paragraph()
p.text = "Lack of Ground Truth: The definition of fraud is ambiguous; unsupervised clustering naturally isolates statistical manipulations without needing perfect labels."
p.level = 0

# Slide 2: Token Selection Criteria
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Mid-Cap Token Selection"
tf = slide.shapes.placeholders[1].text_frame
p = tf.add_paragraph()
p.text = "Selection Criteria for our 5,000 Ethereum Tokens:"
p.level = 0
p = tf.add_paragraph()
p.text = "Instead of blindly trusting unverified 'Market Caps' which are easily manipulated by scammers minting trillions of illiquid tokens, we use activity-based filtering."
p.level = 1
p = tf.add_paragraph()
p.text = "Criteria 1: Top 5,000 tokens by Decentralized Exchange (DEX) trading volume on Uniswap over the last 12 months."
p.level = 1
p = tf.add_paragraph()
p.text = "Criteria 2: Must have a minimum number of unique swaps and active liquidity pools to filter out completely dead tokens."
p.level = 1

# Slide 3: Ground Truth Fraud Labels
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Ground Truth Fraud Baselines"
tf = slide.shapes.placeholders[1].text_frame
p = tf.add_paragraph()
p.text = "We will use these known rug-pulls to validate our unsupervised clusters:"
p.level = 0
p = tf.add_paragraph()
p.text = "Squid Game (SQUID): Crashed precisely on November 1, 2021 (dropped from $2,861 to near $0 in 15 mins)."
p.level = 1
p = tf.add_paragraph()
p.text = "SaveTheKids ($KIDS): Collapsed immediately on launch on June 5, 2021."
p.level = 1
p = tf.add_paragraph()
p.text = "SafeMoon (SFM): Initial massive liquidity crisis on April 20, 2021."
p.level = 1
p = tf.add_paragraph()
p.text = "AnubisDAO (ANKH): $60M stolen during launch on October 29, 2021."
p.level = 1

# Slide 4: Price & Volume Features (Gillian)
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Price & Volume Feature Design"
tf = slide.shapes.placeholders[1].text_frame
p = tf.add_paragraph()
p.text = "Design Choice: STATIC Features for Unsupervised Learning"
p.level = 0
p = tf.add_paragraph()
p.text = "Rationale: Standard clustering algorithms (Isolation Forest, DBSCAN) require flat, fixed-length vectors (1 row per token)."
p.level = 1
p = tf.add_paragraph()
p.text = "Teammate Consistency: All teammates must aggregate their dynamic time-series data into consistent static columns."
p.level = 1
p = tf.add_paragraph()
p.text = "Example Static Features:"
p.level = 1
p = tf.add_paragraph()
p.text = "volume_spike_ratio: Max single-day volume / average daily volume."
p.level = 2
p = tf.add_paragraph()
p.text = "total_days_traded: Aggregated count to identify suspicious bursts."
p.level = 2

# Slide 5: Dune Analytics SQL Queries
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Data Extraction (Dune SQL)"
tf = slide.shapes.placeholders[1].text_frame
p = tf.add_paragraph()
p.text = "Querying dex.trades for Price/Volume Dynamics:"
p.level = 0
p = tf.add_paragraph()
p.text = "Using the Trino (DuneSQL) engine to pull swap volumes and approximate prices across decentralized exchanges."
p.level = 1
p = tf.add_paragraph()
p.text = "The query aggregates by token_a_address and day, computing SUM(amount_usd)."
p.level = 1
p = tf.add_paragraph()
p.text = "Then, a parent CTE calculates the MAX() and AVG() to derive our static volume_spike_ratio and maximum dominance."
p.level = 1
p = tf.add_paragraph()
p.text = "See `price_volume_features.sql` in our repository for the full code."
p.level = 1

# Slide 6: Dune Paid Subscription Benefits
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Dune Plus/Premium Benefits"
tf = slide.shapes.placeholders[1].text_frame
p = tf.add_paragraph()
p.text = "Why we should upgrade the Dune account:"
p.level = 0
p = tf.add_paragraph()
p.text = "Export Limits: Free tier limits UI/API exports, making it hard to download 5,000 tokens' worth of raw data. Paid allows massive programmatic exporting."
p.level = 1
p = tf.add_paragraph()
p.text = "Execution Timeouts: Complex aggregations on `dex.trades` for 5,000 tokens over 12 months will likely time out on free tier compute. Paid tier gets priority execution."
p.level = 1
p = tf.add_paragraph()
p.text = "Privacy: Paid tier allows private queries, protecting our methodology before final evaluation."
p.level = 1

prs.save('Fraud_Detection_Update.pptx')
print("Successfully generated Fraud_Detection_Update.pptx")
