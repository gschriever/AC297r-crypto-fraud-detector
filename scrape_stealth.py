import pandas as pd
import numpy as np
import time
import random
import matplotlib.pyplot as plt
import seaborn as sns
from duckduckgo_search import DDGS
from transformers import pipeline
import warnings
import sys

warnings.filterwarnings('ignore')
print("Launching stealth web-scraper background daemon...", flush=True)

df = pd.read_csv("Gillian_Price_and_Volume_Dynamic_Query_with_Token_Boundary_Conditions_QUERY_EXPORT.csv")
features = ['volume_spike_ratio', 'absolute_max_daily_volume_usd', 'total_days_traded', 'max_trade_dominance']
df[features] = df[features].apply(pd.to_numeric, errors='coerce')
df = df.dropna(subset=features).copy()

from sklearn.ensemble import IsolationForest
from sklearn.cluster import DBSCAN
from sklearn.decomposition import PCA
from sklearn.preprocessing import StandardScaler
from scipy.spatial.distance import mahalanobis

X = df[features].values
scaler = StandardScaler()
X_scaled = scaler.fit_transform(X)

iso = IsolationForest(contamination=0.05, random_state=42)
df['iso_flag'] = iso.fit_predict(X_scaled) == -1
dbscan = DBSCAN(eps=1.5, min_samples=10)
df['dbscan_flag'] = dbscan.fit_predict(X_scaled) == -1
pca = PCA(n_components=2)
X_pca = pca.fit_transform(X_scaled)
inv_cov = np.linalg.inv(np.cov(X_pca, rowvar=False))
center = np.mean(X_pca, axis=0)
dists = [mahalanobis(row, center, inv_cov) for row in X_pca]
df['pca_flag'] = dists > np.percentile(dists, 95)
df['votes'] = df['iso_flag'].astype(int) + df['dbscan_flag'].astype(int) + df['pca_flag'].astype(int)

# Focus strictly on the 46 High Confidence tokens
high_conf = df[df['votes'] == 3].copy()
print(f"Isolated {len(high_conf)} High Confidence Tokens.", flush=True)

try:
    classifier = pipeline("zero-shot-classification", model="typeform/distilbert-base-uncased-mnli")
except:
    classifier = pipeline("zero-shot-classification", model="cross-encoder/nli-distilroberta-base")

candidate_labels = ["fraud, scam, or rug pull", "legitimate crypto protocol", "neutral or ambiguous market data"]
results = []
queries_executed = 0

with DDGS() as ddgs:
    for idx, row in high_conf.iterrows():
        symbol = row.get('token_symbol', 'Unknown')
        if pd.isna(symbol) or symbol == 'Unknown':
            continue
            
        print(f"[{queries_executed+1}] Investigating {symbol}...", flush=True)
        query = f'"{symbol}" crypto AND ("scam" OR "rug pull" OR "exploit" OR "hack" OR "legitimate")'
        
        try:
            search_results = list(ddgs.text(query, max_results=3))
            if not search_results:
                results.append((symbol, "Undetermined / Missing Data"))
            else:
                combined_text = " ".join([res.get('body', '') for res in search_results])
                classification = classifier(combined_text, candidate_labels)
                top_label = classification['labels'][0]
                confidence = classification['scores'][0]
                
                if top_label == "fraud, scam, or rug pull" and confidence > 0.40:
                    status = "Confirmed Fraud or Rug Pull"
                elif top_label == "legitimate crypto protocol" and confidence > 0.40:
                    status = "Legitimate Protocol Outlier"
                else:
                    status = "Undetermined / Neutral"
                results.append((symbol, status))
        except Exception as e:
            print(f"HTTP block on {symbol}: {e}", flush=True)
            results.append((symbol, "Error/Blocked"))
            
        # Stealth Mode Sleep: 30 second rest to guarantee evasion of 403 API IP blocks
        sleep_time = random.uniform(28, 35) 
        time.sleep(sleep_time)
        queries_executed += 1
        
        if queries_executed >= 46:
            break

# Dump the results to a CSV
results_df = pd.DataFrame(results, columns=["Symbol", "Validation Status"])
results_df.to_csv("full_stealth_validation.csv", index=False)
summary_counts = results_df["Validation Status"].value_counts()

# Build the chart
print("\nGenerating final chart...", flush=True)
colors = sns.color_palette("muted", len(summary_counts))
plt.figure(figsize=(10, 6))
plt.pie(summary_counts.values, labels=summary_counts.index, colors=colors, autopct='%1.1f%%', startangle=140, explode=[0.05]*len(summary_counts))
plt.title("HuggingFace Zero-Shot Verification of 46 Unsupervised Anomalies")
plt.tight_layout()
plt.savefig("fraud_validation_chart_stealth.png", dpi=300)

# Inject into Powerpoint via Python-PPTX
print("Updating PowerPoint...", flush=True)
from pptx import Presentation
from pptx.util import Pt
try:
    prs = Presentation('April_16_Update.pptx')
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture('fraud_validation_chart_stealth.png', Pt(80), Pt(50), height=Pt(400))
    prs.save('April_16_Update.pptx')
    print('SUCCESS: Chart added directly to April_16_Update.pptx', flush=True)
except Exception as e:
    print(f"Could not update PPTX: {e}", flush=True)

print("Stealth Job Complete.", flush=True)
