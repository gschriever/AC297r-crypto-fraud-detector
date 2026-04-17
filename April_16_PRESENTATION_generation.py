import collections
import collections.abc
from pptx import Presentation
from pptx.util import Pt

prs = Presentation()

# Slide 1: Title
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
slide.shapes.title.text = "Unsupervised Fraud Detection Validation\nFinal Results"
slide.placeholders[1].text = "AC297r\nGillian Schriever | Jiahui (Cecilia) Cai | Zhilin Chen | Jinghan Huang"

# Slide 2: Evaluation Methodology
bullet_slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(bullet_slide_layout)
slide.shapes.title.text = "Evaluation Methodology"
tf = slide.placeholders[1].text_frame
tf.text = "Goal: Isolate fraudulent behavior natively without labeled training data."
p = tf.add_paragraph()
p.text = "Evaluated 4,334 mid-cap ERC-20 tokens via Dune Analytics."
p = tf.add_paragraph()
p.text = "Deployed 3 Independent Unsupervised Models:"
p = tf.add_paragraph()
p.text = "1. Isolation Forest (Detects sparse outliers)"
p.level = 1
p = tf.add_paragraph()
p.text = "2. DBSCAN (Detects density anomalies)"
p.level = 1
p = tf.add_paragraph()
p.text = "3. PCA + Mahalanobis Distance (Detects multi-dimensional extremes)"
p.level = 1

notes_slide = slide.notes_slide
notes_slide.notes_text_frame.text = ("TALKING SCRIPT: Our primary challenge was proving that an algorithm could detect a scam without ever being told what a scam looks like. We took the 4,334 active tokens extracted from our Dune SQL pipeline and passed them through three distinct unsupervised models. We used Isolation Forest for sparse outlier detection, DBSCAN to find tokens sitting outside of normal behavioral clusters, and PCA paired with Mahalanobis distance to measure multi-dimensional extremes. Our hypothesis was simple: while one model might occasionally flag a false positive, any token flagged by all three models simultaneously would mathematically represent unnatural, extreme trading behavior indicative of manipulation.")

# Slide 3: The Consensus Approach
slide = prs.slides.add_slide(bullet_slide_layout)
slide.shapes.title.text = "The Consensus Approach"
tf = slide.placeholders[1].text_frame
tf.text = "To eliminate false positives, we established a strict voting threshold."
p = tf.add_paragraph()
p.text = "Tokens flagged by at least 2 models: 147"
p = tf.add_paragraph()
p.text = "Tokens flagged by ALL 3 models (High Confidence): 46"
p = tf.add_paragraph()
p.text = "The Top 46 tokens represent the most mathematically extreme ~1% of the entire Ethereum mid-cap market."

notes_slide = slide.notes_slide
notes_slide.notes_text_frame.text = ("TALKING SCRIPT: When we ran the models, we wanted to ensure the highest possible confidence before accusing a token of manipulation. We established a strict consensus protocol. Out of 4,334 tokens, 147 triggered at least two models. More importantly, only 46 tokens managed to trigger all three anomaly detectors at the exact same time. By isolating this 1% of the market as our 'High Confidence' tier, our final step was to manually investigate these extreme outliers to see if the math correlated with reality.")

# Slide 4: Ground Truth Validation
slide = prs.slides.add_slide(bullet_slide_layout)
slide.shapes.title.text = "Ground Truth Validation"
tf = slide.placeholders[1].text_frame
tf.text = "The model successfully isolated highly documented, real-world scams:"
p = tf.add_paragraph()
p.text = "YODA (Yoda Coin): Confirmed $129k Rug Pull."
p.level = 1
p = tf.add_paragraph()
p.text = "SURGE (xSURGE): Confirmed Exit Scam."
p.level = 1
p = tf.add_paragraph()
p.text = "WLFI (World Liberty Financial): Mainstream Pump-and-Dump Controversies."
p.level = 1
p = tf.add_paragraph()
p.text = "DFED: Confirmed Honeypot Exploit."
p.level = 1

notes_slide = slide.notes_slide
notes_slide.notes_text_frame.text = ("TALKING SCRIPT: This is our ultimate validation slide and arguably the most crucial finding of the capstone. We looked up the top anomalies flagged by the algorithm—tokens like YODA, SURGE, WLFI, and DFED. Astonishingly, our web research confirmed that these tokens were not just statistical flukes; they are famous, heavily-documented scams. YODA's creators abruptly deleted their social accounts and drained $129k in a classic rug pull. SURGE suffered a massive exit scam. WLFI faced intense mainstream allegations of extreme insider manipulation and massive whale dumping. DFED is explicitly flagged by blockchain security scanners as a honeypot (meaning buyers are permanently blocked from selling). We successfully proved that without a single line of training data, our feature engineering captured the physical signature of fraud on the blockchain.")

prs.save('Final_Validation_Presentation.pptx')
print("Saved Final_Validation_Presentation.pptx")
