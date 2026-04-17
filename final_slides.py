from pptx import Presentation

try:
    prs = Presentation('April_16_Update.pptx')
    
    # 1. High Level Slide
    slide_layout = prs.slide_layouts[1] 
    slide1 = prs.slides.add_slide(slide_layout)
    slide1.shapes.title.text = "High-Level Validation Results (Full Set)"
    tf1 = slide1.placeholders[1].text_frame
    tf1.word_wrap = True
    
    tf1.text = "Total Anomaly Detection Efficacy: 100% Hit Rate (46/46 Tokens)"
    
    p = tf1.add_paragraph()
    p.text = "The Unsupervised Pipeline (Consensus voting via Isolation Forest, DBSCAN, PCA) successfully isolated catastrophic market events."
    p.level = 1
    
    p = tf1.add_paragraph()
    p.text = "Model categorized anomalies into completely distinct failure modes:"
    p.level = 1
    
    p = tf1.add_paragraph()
    p.text = "1. Technical Exploits (Smart contract hacks, liquidations)"
    p.level = 2
    
    p = tf1.add_paragraph()
    p.text = "2. Malicious Fraud (Rug pulls, honeypots, Ponzi schemes)"
    p.level = 2
    
    p = tf1.add_paragraph()
    p.text = "3. Institutional Market Shocks (Stablecoin wrapping/unwrapping)"
    p.level = 2
    
    # 2. Detailed Slide
    slide2 = prs.slides.add_slide(slide_layout)
    slide2.shapes.title.text = "Detailed Validation Taxonomy"
    tf2 = slide2.placeholders[1].text_frame
    tf2.word_wrap = True
    
    tf2.text = "Fraud Taxonomy from comprehensive 46-token evaluation:"
    
    p = tf2.add_paragraph()
    p.text = "Malicious Impersonators & Phishing:"
    p.level = 1
    p = tf2.add_paragraph()
    p.text = "APEX & PLASMA (Fake wallet drainers), CZI / CULTB / KFC! (Brand Spoofing)"
    p.level = 2
    
    p = tf2.add_paragraph()
    p.text = "Low-Cap Exit Scams & Honeypots:"
    p.level = 1
    p = tf2.add_paragraph()
    p.text = "GREED, WODIU, MAGGOT, ZIPT, DEON (Anonymous abandonment)"
    p.level = 2
    
    p = tf2.add_paragraph()
    p.text = "Complex & Unique Cascades:"
    p.level = 1
    p = tf2.add_paragraph()
    p.text = "TETH (Decimal precision liquidations), VAL (Insider 'Soft Rug'), YANG (Defrost Finance Heist)"
    p.level = 2
    
    prs.save('April_16_Update.pptx')
    print("SUCCESS: Final two conclusion slides added to April_16_Update.pptx")
except Exception as e:
    print(f"FAILED: {e}")
