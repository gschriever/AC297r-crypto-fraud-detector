from pptx import Presentation

try:
    prs = Presentation('April_17_Update.pptx')
    slide_layout = prs.slide_layouts[1] 
    slide1 = prs.slides.add_slide(slide_layout)
    slide1.shapes.title.text = "Ground Truth Strict Validation (True vs False Positives)"
    tf1 = slide1.placeholders[1].text_frame
    tf1.word_wrap = True
    
    tf1.text = "Defining 'Scam' strictly as malicious developer fraud (rug pulls, pump & dumps):"
    
    p = tf1.add_paragraph()
    p.text = "True Positives (~80%): Intentional Developer Fraud"
    p.level = 1
    
    p = tf1.add_paragraph()
    p.text = "Massive rug pulls (TARO, GREED), honeypots (LUA), corporate spoofing (CZI, CULTB), and massive phishing drainers (PLASMA, APEX)."
    p.level = 2

    p = tf1.add_paragraph()
    p.text = "False Positives (~20%): Benign Mechanics & Market Shocks"
    p.level = 1
    
    p = tf1.add_paragraph()
    p.text = "Institutional Liquidity: 3Crv, USDP, USDL, aEthWETH, aEthWBTC (Whale / institutional rebalancing)."
    p.level = 2
    
    p = tf1.add_paragraph()
    p.text = "Contagion & Cascades: MIM (UST De-pegging Contagion), CORE (Algorithmic liquidation cascade)."
    p.level = 2
    
    p = tf1.add_paragraph()
    p.text = "Bugs & External Hacks: TETH (Decimal mismatch bug), SLP (Ronin $622M Hack), BNT (Bancor external exploit protocol freeze)."
    p.level = 2

    prs.save('April_17_Update_v2.pptx')
    print("SUCCESS")
except Exception as e:
    print(f"FAILED: {e}")
