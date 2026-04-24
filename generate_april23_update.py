from pptx import Presentation

QUESTIONS = [
    "Headline framing — continuous score vs. discrete consensus voting",
    "Calibration sample size",
    "Log-transform on skewed features",
    "Calibrator extrapolation handling",
    "DBSCAN's role in the ensemble",
]

UPDATES = [
    "We stopped throwing away data",
    "We switched from yes/no voting to a dimmer switch",
    'We added a second layer that estimates "how likely is this suspicious behavior to actually be fraud?"',
    "We added a way to explain why the model flagged each token",
    "We ran a stress test on our thresholds",
]


def add_bullet_slide(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.text = bullets[0]
    for b in bullets[1:]:
        p = tf.add_paragraph()
        p.text = b


def main():
    prs = Presentation()

    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = "April 23 Update"
    title_slide.placeholders[1].text = (
        "AC297r — Unsupervised Crypto Fraud Detector\n"
        "Gillian Schriever | Jiahui (Cecilia) Cai | Zhilin Chen | Jinghan Huang"
    )

    add_bullet_slide(prs, "Top 5 Questions for Our Professor", QUESTIONS)
    add_bullet_slide(prs, "Top 5 Updates This Week", UPDATES)

    out = "April_23_Update.pptx"
    prs.save(out)
    print(f"Saved {out}")


if __name__ == "__main__":
    main()
