from pptx import Presentation

try:
    prs = Presentation('April_17_Update.pptx')
    slide_layout = prs.slide_layouts[1] 
    slide1 = prs.slides.add_slide(slide_layout)
    slide1.shapes.title.text = "Severity Scaling via Consensus"
    tf1 = slide1.placeholders[1].text_frame
    tf1.word_wrap = True
    
    tf1.text = "Key Insight: The severity of the market event scales perfectly with the number of models flagging it."
    
    p = tf1.add_paragraph()
    p.text = "3-Vote Consensus (Isolation Forest + DBSCAN + PCA):"
    p.level = 1
    
    p = tf1.add_paragraph()
    p.text = "Catastrophic systemic failures and massive protocol hacks (e.g., $600M Ronin Hack, 50% CORE liquidations)."
    p.level = 2

    p = tf1.add_paragraph()
    p.text = "Institutional volume shocks (Stablecoin mint/burn cascades)."
    p.level = 2
    
    p = tf1.add_paragraph()
    p.text = "1-Vote Anomalies (Single Model Flag):"
    p.level = 1
    
    p = tf1.add_paragraph()
    p.text = "Basic, low-liquidity meme coin rug pulls (e.g., RATO, $ANIBTC)."
    p.level = 2
    
    p = tf1.add_paragraph()
    p.text = "Name collisions targeting slang and hype keywords (e.g., FOMO, MIS)."
    p.level = 2

    prs.save('April_17_Update.pptx')
    print("SUCCESS: Slide added to April_17_Update.pptx")
except Exception as e:
    print(f"FAILED: {e}")
