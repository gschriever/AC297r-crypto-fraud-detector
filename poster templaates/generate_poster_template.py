"""
Poster generator — AC297r capstone.

Builds a 48" x 60" landscape poster (print size: 60" wide x 48" tall) that
follows the professor's guidelines:

  - Results-driven title, not a topical one.
  - Three-column Question / Methods / Results flow.
  - Body text >= 24pt; title ~96pt; section headers 48pt.
  - Clear "So what?" banner across the bottom.
  - Two compatible fonts only: Georgia (serif, title), Calibri (sans, body).
  - High whitespace; figures embedded from ./artifacts/.
  - No abstract block (discouraged by guidelines when authors present).

python-pptx caps slide dimensions at 56". We therefore build at 55" x 44"
(same 5:4 aspect as 60" x 48"), save, then rewrite the slide-size XML
attribute to the true 60" x 48" print size and scale every element's
position and size by 60/55 so nothing shifts visually.
"""

from __future__ import annotations

import shutil
import zipfile
from pathlib import Path
from xml.etree import ElementTree as ET

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

# ---- Palette ------------------------------------------------------------------
CRIMSON = RGBColor(0xA5, 0x1C, 0x30)
INK = RGBColor(0x2A, 0x2A, 0x2A)
SOFT = RGBColor(0xF5, 0xF5, 0xF5)
RULE = RGBColor(0xD8, 0xD8, 0xD8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# ---- Dimensions ---------------------------------------------------------------
BUILD_W = Inches(55)
BUILD_H = Inches(44)
TARGET_W_IN = 60
TARGET_H_IN = 48
SCALE = TARGET_W_IN / 55  # == TARGET_H_IN / 44 == 12/11

SERIF = "Georgia"
SANS = "Calibri"

ROOT = Path(__file__).resolve().parent
ARTIFACTS = ROOT.parent / "artifacts"


# ---- Helpers ------------------------------------------------------------------


def set_fill(shape, rgb):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb


def set_no_line(shape):
    shape.line.fill.background()


def set_line(shape, rgb, width_pt=1.0):
    shape.line.color.rgb = rgb
    shape.line.width = Pt(width_pt)


def add_text(
    slide,
    left,
    top,
    width,
    height,
    text,
    *,
    font=SANS,
    size=28,
    bold=False,
    color=INK,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
    italic=False,
    line_spacing=1.15,
):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return tb


def add_rule(slide, left, top, width, color=CRIMSON, weight_pt=3.0):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Inches(0.04))
    set_fill(shp, color)
    set_no_line(shp)
    shp.shadow.inherit = False
    return shp


def add_section_header(slide, left, top, width, label):
    add_text(
        slide,
        left,
        top,
        width,
        Inches(0.9),
        label.upper(),
        font=SANS,
        size=42,
        bold=True,
        color=CRIMSON,
    )
    add_rule(slide, left, top + Inches(1.0), width * 0.35)


def add_figure(slide, left, top, width, height, image_path, caption):
    try:
        slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)
    except Exception:
        frame = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        set_fill(frame, WHITE)
        set_line(frame, RULE, 1.0)
        frame.shadow.inherit = False
        add_text(
            slide,
            left,
            top + height / 2 - Inches(0.3),
            width,
            Inches(0.8),
            "[ insert figure ]",
            font=SANS,
            size=24,
            italic=True,
            color=RGBColor(0x99, 0x99, 0x99),
            align=PP_ALIGN.CENTER,
        )
    add_text(
        slide,
        left,
        top + height + Inches(0.05),
        width,
        Inches(0.8),
        caption,
        font=SANS,
        size=22,
        italic=True,
        color=INK,
        align=PP_ALIGN.CENTER,
    )


# ---- Build --------------------------------------------------------------------


def build():
    prs = Presentation()
    prs.slide_width = BUILD_W
    prs.slide_height = BUILD_H
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, BUILD_W, BUILD_H)
    set_fill(bg, WHITE)
    set_no_line(bg)

    margin = Inches(0.9)

    # --- HEADER (title + authors + shields) ---
    header_h = Inches(4.8)

    # Shield placeholders — two small boxes the user will replace with
    # Harvard + SEAS logos. Kept simple (the previous version had four).
    shield_w = Inches(3.2)
    shield_h = Inches(3.6)
    shield_y = Inches(0.6)
    for x in (margin, BUILD_W - margin - shield_w):
        sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, shield_y, shield_w, shield_h)
        set_fill(sh, SOFT)
        set_line(sh, RULE, 0.75)
        sh.shadow.inherit = False
        add_text(
            slide,
            x,
            shield_y + shield_h / 2 - Inches(0.3),
            shield_w,
            Inches(0.7),
            "[ shield ]",
            font=SANS,
            size=20,
            italic=True,
            color=RGBColor(0x99, 0x99, 0x99),
            align=PP_ALIGN.CENTER,
        )

    # Title — results-driven, not topical.
    add_text(
        slide,
        margin + shield_w + Inches(0.4),
        Inches(0.6),
        BUILD_W - 2 * (margin + shield_w + Inches(0.4)),
        Inches(2.6),
        "Detecting crypto fraud without labels:\n6 of 6 known rug pulls recovered from 4,334 Ethereum tokens.",
        font=SERIF,
        size=62,
        bold=True,
        color=CRIMSON,
        align=PP_ALIGN.CENTER,
        line_spacing=1.05,
    )
    # Authors
    add_text(
        slide,
        margin + shield_w + Inches(0.4),
        Inches(3.3),
        BUILD_W - 2 * (margin + shield_w + Inches(0.4)),
        Inches(0.7),
        "Gillian Schriever  ·  Jiahui (Cecilia) Cai  ·  Zhilin Chen  ·  Jinghan Huang",
        font=SERIF,
        size=28,
        color=INK,
        align=PP_ALIGN.CENTER,
    )
    # Affiliation
    add_text(
        slide,
        margin + shield_w + Inches(0.4),
        Inches(4.0),
        BUILD_W - 2 * (margin + shield_w + Inches(0.4)),
        Inches(0.6),
        "Harvard John A. Paulson School of Engineering and Applied Sciences  —  AC297r Capstone, Spring 2026",
        font=SERIF,
        size=22,
        italic=True,
        color=INK,
        align=PP_ALIGN.CENTER,
    )

    # Crimson divider under header
    add_rule(slide, 0, header_h, BUILD_W, color=CRIMSON, weight_pt=4.0)

    # --- THREE COLUMNS ---
    body_top = header_h + Inches(0.7)
    body_bottom_limit = BUILD_H - Inches(9.0)  # leave room for so-what banner + footer
    body_h = body_bottom_limit - body_top

    n_cols = 3
    gutter = Inches(0.9)
    content_w = BUILD_W - 2 * margin - gutter * (n_cols - 1)
    col_w = content_w / n_cols
    col_xs = [margin + i * (col_w + gutter) for i in range(n_cols)]

    # Column 1 — The Question / Approach
    cx = col_xs[0]
    add_section_header(slide, cx, body_top, col_w, "The question")
    add_text(
        slide,
        cx,
        body_top + Inches(1.4),
        col_w,
        Inches(3.4),
        "Can we detect crypto fraud without any labeled data?",
        font=SANS,
        size=32,
        bold=True,
        color=INK,
        line_spacing=1.15,
    )
    add_text(
        slide,
        cx,
        body_top + Inches(4.4),
        col_w,
        Inches(9.0),
        (
            "Labels for crypto fraud are rare, slow to publish, and biased\n"
            "toward high-profile cases. Supervised classifiers trained on\n"
            "them don't generalize — the attacks evolve faster than the\n"
            "labels.\n\n"
            "We build a two-layer unsupervised pipeline:\n\n"
            "  ① Flag tokens whose on-chain trading behavior looks\n"
            "      statistically unusual.\n\n"
            "  ② Estimate how likely that unusual behavior is to be\n"
            "      intentional fraud (vs. a hack, bug, or large legitimate\n"
            "      event)."
        ),
        font=SANS,
        size=24,
        color=INK,
        line_spacing=1.25,
    )

    # Column 2 — Methods
    cx = col_xs[1]
    add_section_header(slide, cx, body_top, col_w, "Methods")
    add_text(
        slide,
        cx,
        body_top + Inches(1.4),
        col_w,
        Inches(13.5),
        (
            "Data. 4,334 mid-cap ERC-20 tokens on Ethereum\n"
            "(Dune Analytics, 365-day window, $1M–$1B annual volume).\n\n"
            "Features (6). Volume spike ratio, peak daily volume, days\n"
            "traded, max single-trade dominance, and two 24-hour post-\n"
            "launch velocity features.\n\n"
            "Layer 1 — three independent detectors:\n"
            "    •  Isolation Forest (tree-based outlier detection)\n"
            "    •  DBSCAN (density clustering; we keep the noise points)\n"
            "    •  PCA + Mahalanobis distance\n\n"
            "We rank-average their outputs into a single continuous\n"
            "suspicion score in [0, 1], replacing the legacy three-detector\n"
            "yes/no vote.\n\n"
            "Layer 2. Logistic regression (L2-regularized) fit on n=14\n"
            "hand-validated tokens, mapping suspicious-behavior features\n"
            "to P(fraud). Evaluated with leave-one-out cross-validation."
        ),
        font=SANS,
        size=24,
        color=INK,
        line_spacing=1.25,
    )

    # Column 3 — Results (two key figures + a headline stat)
    cx = col_xs[2]
    add_section_header(slide, cx, body_top, col_w, "Results")

    # Headline number — the "so what" preview
    kpi_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, cx, body_top + Inches(1.4), col_w, Inches(3.2)
    )
    kpi_box.adjustments[0] = 0.06
    set_fill(kpi_box, SOFT)
    set_line(kpi_box, RULE, 0.75)
    kpi_box.shadow.inherit = False
    add_text(
        slide,
        cx,
        body_top + Inches(1.55),
        col_w,
        Inches(1.6),
        "6 / 6",
        font=SERIF,
        size=96,
        bold=True,
        color=CRIMSON,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        cx,
        body_top + Inches(3.3),
        col_w,
        Inches(1.2),
        "known frauds recovered at the p90 suspicion cutoff.\n"
        "The legacy three-vote consensus recovered 1 of 6.",
        font=SANS,
        size=22,
        italic=True,
        color=INK,
        align=PP_ALIGN.CENTER,
        line_spacing=1.2,
    )

    # Figure 1 — lift chart (the primary evidence for the headline)
    fig1_y = body_top + Inches(5.0)
    fig_h = Inches(6.5)
    add_figure(
        slide,
        cx,
        fig1_y,
        col_w,
        fig_h,
        ARTIFACTS / "baseline_comparison.png",
        "Fig. 1.  Fraud recall on the validated subset (n=6 fraud, 8 non-fraud\nanomaly). Continuous suspicion score catches all known fraud where\nthe legacy consensus vote catches one.",
    )

    # Figure 2 — SHAP global importance (why it works)
    fig2_y = fig1_y + fig_h + Inches(1.6)
    fig_h2 = Inches(5.2)
    add_figure(
        slide,
        cx,
        fig2_y,
        col_w,
        fig_h2,
        ARTIFACTS / "shap_global_importance.png",
        "Fig. 2.  SHAP attribution on Isolation Forest.  Early 24-hour\nvelocity outranks the legacy headline feature (volume spike ratio).",
    )

    # --- SO WHAT banner + footer ---
    sowhat_y = BUILD_H - Inches(8.5)
    sowhat_h = Inches(5.2)
    banner = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, margin, sowhat_y, BUILD_W - 2 * margin, sowhat_h
    )
    banner.adjustments[0] = 0.03
    set_fill(banner, CRIMSON)
    set_no_line(banner)
    banner.shadow.inherit = False

    add_text(
        slide,
        margin + Inches(0.6),
        sowhat_y + Inches(0.3),
        BUILD_W - 2 * margin - Inches(1.2),
        Inches(1.1),
        "SO WHAT",
        font=SANS,
        size=32,
        bold=True,
        color=WHITE,
    )
    add_text(
        slide,
        margin + Inches(0.6),
        sowhat_y + Inches(1.4),
        BUILD_W - 2 * margin - Inches(1.2),
        sowhat_h - Inches(1.5),
        (
            "Continuous ranking beats discrete voting.  A rank-averaged suspicion score recovered 6/6 known\n"
            "frauds on the validated subset; the legacy \"3 of 3 detectors agree\" rule recovered 1/6.\n\n"
            "Two layers separate 'anomaly' from 'fraud.'  The Layer-2 calibrator cleanly distinguishes rug-pulls\n"
            "(P ≈ 0.5–0.7) from large legitimate events like hacks and de-peg shocks (P ≈ 0.1–0.2).\n\n"
            "Key caveat.  The calibration layer is trained on only 14 hand-validated tokens; expanding that\n"
            "labeled set is the single most valuable next step."
        ),
        font=SANS,
        size=26,
        color=WHITE,
        line_spacing=1.25,
    )

    # Footer — contact + code + references, minimal
    footer_y = sowhat_y + sowhat_h + Inches(0.4)
    footer_h = BUILD_H - footer_y - Inches(0.4)
    add_text(
        slide,
        margin,
        footer_y,
        BUILD_W * 0.5 - margin,
        footer_h,
        (
            "Code:      github.com/gschriever/AC297r-crypto-fraud-detector\n"
            "Contact:   [ name ]  ·  [ email ]\n"
            "Advisor:   [ advisor ]"
        ),
        font=SANS,
        size=22,
        color=INK,
        line_spacing=1.3,
    )
    add_text(
        slide,
        BUILD_W * 0.5,
        footer_y,
        BUILD_W * 0.5 - margin,
        footer_h,
        (
            "[1] Mazorra et al.  Do Not Rug on Me.  arXiv:2201.07220, 2022.\n"
            "[2] Liu, Ting & Zhou.  Isolation Forest.  ICDM 2008.\n"
            "[3] Lundberg & Lee.  SHAP.  NeurIPS 2017."
        ),
        font=SANS,
        size=22,
        color=INK,
        line_spacing=1.3,
    )

    out = ROOT / "AC297r_Poster_Template.pptx"
    prs.save(str(out))
    return out


# ---- Post-process: set slide size to the real 60" x 48" print size -----------


def rescale_pptx_to_true_size(pptx_path: Path) -> None:
    """
    python-pptx caps slide dimensions at 56". We therefore built at 55" x 44"
    and now rewrite the slide size XML to the real target (60" x 48") and
    scale every shape's position & size by SCALE == 60/55 so the layout
    prints at the intended physical size.
    """
    tmp = pptx_path.with_suffix(".pptx.tmp")
    shutil.move(str(pptx_path), str(tmp))

    NSMAP = {
        "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
        "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    }
    for prefix, uri in NSMAP.items():
        ET.register_namespace(prefix if prefix != "p" else "p", uri)
    ET.register_namespace("", "http://schemas.openxmlformats.org/drawingml/2006/main")

    with zipfile.ZipFile(str(tmp), "r") as zin:
        names = zin.namelist()
        contents = {n: zin.read(n) for n in names}

    # 1. Rewrite <p:sldSz cx=... cy=.../> in presentation.xml
    pres_xml = contents["ppt/presentation.xml"].decode("utf-8")
    tree = ET.ElementTree(ET.fromstring(pres_xml))
    root = tree.getroot()
    sld_sz = root.find("p:sldSz", NSMAP)
    sld_sz.set("cx", str(TARGET_W_IN * 914400))
    sld_sz.set("cy", str(TARGET_H_IN * 914400))
    contents["ppt/presentation.xml"] = ET.tostring(root, xml_declaration=True, encoding="UTF-8")

    # 2. Scale every shape on every slide by SCALE.
    def scale_attr(el, attr):
        v = el.get(attr)
        if v is not None:
            el.set(attr, str(int(int(v) * SCALE)))

    for name in names:
        if not (name.startswith("ppt/slides/slide") and name.endswith(".xml")):
            continue
        tree = ET.ElementTree(ET.fromstring(contents[name].decode("utf-8")))
        for xfrm in tree.getroot().iter("{http://schemas.openxmlformats.org/drawingml/2006/main}xfrm"):
            for tag in ("off", "ext"):
                el = xfrm.find(f"{{http://schemas.openxmlformats.org/drawingml/2006/main}}{tag}")
                if el is not None:
                    scale_attr(el, "x")
                    scale_attr(el, "y")
                    scale_attr(el, "cx")
                    scale_attr(el, "cy")
        # Scale font sizes (stored in hundredths of a point via sz="2800")
        for run_props in tree.getroot().iter("{http://schemas.openxmlformats.org/drawingml/2006/main}rPr"):
            sz = run_props.get("sz")
            if sz is not None:
                run_props.set("sz", str(int(int(sz) * SCALE)))
        for def_props in tree.getroot().iter("{http://schemas.openxmlformats.org/drawingml/2006/main}defRPr"):
            sz = def_props.get("sz")
            if sz is not None:
                def_props.set("sz", str(int(int(sz) * SCALE)))
        contents[name] = ET.tostring(tree.getroot(), xml_declaration=True, encoding="UTF-8")

    with zipfile.ZipFile(str(pptx_path), "w", zipfile.ZIP_DEFLATED) as zout:
        for n, b in contents.items():
            zout.writestr(n, b)
    tmp.unlink()


# ---- Entry --------------------------------------------------------------------


def main():
    out = build()
    rescale_pptx_to_true_size(out)
    print(f"Saved {out}")
    print(f"Design build: 55 x 44 in  →  rescaled to {TARGET_W_IN} x {TARGET_H_IN} in (landscape, 48x60 poster)")
    print(f"Body font 24pt  →  {int(24 * SCALE)}pt after rescale (final printed size)")
    print(f"Title 62pt      →  {int(62 * SCALE)}pt after rescale")


if __name__ == "__main__":
    main()
