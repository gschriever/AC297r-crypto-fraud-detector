import collections
import collections.abc
from pptx import Presentation
from pptx.util import Pt

prs = Presentation()

# Slide 1: Title
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Detecting Pump-and-Dump in Small-Cap Ethereum Tokens\nApril 9 Update"
subtitle.text = "AC297r\nGillian Schriever | Jiahui (Cecilia) Cai | Zhilin Chen | Jinghan Huang"

# Slide 2: Team Updates
bullet_slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(bullet_slide_layout)
title_shape = slide.shapes.title
body_shape = slide.placeholders[1]
title_shape.text = "Team Updates"
tf = body_shape.text_frame
tf.text = "Successfully Secured Dune Analytics Pro Tier Access"
p = tf.add_paragraph()
p.text = "Defined Boundary Conditions: The Universal 'Master Token List'"
p.level = 0
p = tf.add_paragraph()
p.text = "master_token_list AS (SELECT token_address FROM raw_trades ... LIMIT 5000)"
p.level = 1
# Optionally set font size
p.font.size = Pt(14)
p.font.name = "Courier New"
p = tf.add_paragraph()
p.text = "Executed Initial Price/Volume Queries & Exported Raw CSV"
p.level = 0

notes_slide = slide.notes_slide
text_frame = notes_slide.notes_text_frame
text_frame.text = ("Notes on Master CTE: We resolved a core issue where different modules were pulling 100 randomly different tokens. By writing a single logical boundary condition targeting $1M-$1B total volume, we ensured that every teammate queries the exact same 5,000 active mid-caps, guaranteeing our CSVs will join perfectly. \n\nPro Tier Access: Upgrading bypassed all API export ceilings and execution timeouts, allowing us to safely process the massive dex.trades database covering the entire 365-day spectrum.")

# Slide 3: Next Steps
slide = prs.slides.add_slide(bullet_slide_layout)
title_shape = slide.shapes.title
body_shape = slide.placeholders[1]
title_shape.text = "Next Steps"
tf = body_shape.text_frame
tf.text = "1. Execute All Modules & Export CSVs"
p = tf.add_paragraph()
p.text = "2. Python Pandas Data Consolidation"
p.level = 0
p = tf.add_paragraph()
p.text = "3. Deploy Unsupervised Algorithmic Modeling"
p.level = 0
p = tf.add_paragraph()
p.text = "4. Validate Anomalies against Ground Truth Baselines"
p.level = 0
p = tf.add_paragraph()
p.text = "5. (Time Permitting) Build Interactive Streamlit Dashboard"
p.level = 0

notes_slide = slide.notes_slide
text_frame = notes_slide.notes_text_frame
text_frame.text = ("1. All four teammates need to run their completed Dune logic (incorporating the Master CTE list) and export their 5,000-row features. \n\n2. We will programmatically clean and merge the CSVs matching on token_address. \n\n3. Push the master spreadsheet into the clustering algorithm to calculate anomaly scores. \n\n4. Check if known scams (e.g., AnubisDAO, Mutant Ape Planet) are mathematically captured in the anomaly clusters. \n\n5. If time allows, wrap the findings in a dashboard to visualize exactly WHY a specific token was flagged.")

# Slide 4: Open Questions
slide = prs.slides.add_slide(bullet_slide_layout)
title_shape = slide.shapes.title
body_shape = slide.placeholders[1]
title_shape.text = "Open Questions & Design Explorations"
tf = body_shape.text_frame
tf.text = "Algorithmic Approach: Isolation Forest vs. DBSCAN?"
p = tf.add_paragraph()
p.text = "Dashboard Design: What metrics provide the most value for frontend visualization?"
p.level = 0

notes_slide = slide.notes_slide
text_frame = notes_slide.notes_text_frame
text_frame.text = ("Algorithm Considerations: Our proposal considered standard anomaly detection. Suggestion: We should lean heavily into Isolation Forests. They are highly efficient and effectively handle high-dimensional tabular data without requiring massive pre-scaling assumptions, compared to DBSCAN which struggles with varying-density data clouds.\n\nDashboard Suggestions: If we build the dashboard, should we display the feature breakdown for specific suspicious tokens? Suggestion: We should highlight the token's specific anomaly score compared to the 'market average' (e.g., showing a user exactly how extreme a token's volume spike ratio actually was compared to normal mid-caps).")

prs.save('April_9_Update.pptx')
print("Successfully saved April_9_Update.pptx")
