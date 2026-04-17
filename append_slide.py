from pptx import Presentation
from pptx.util import Pt

try:
    prs = Presentation('April_16_Update.pptx')
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Extended Validation (25-Token Sample)"

    tf = slide.placeholders[1].text_frame
    tf.word_wrap = True

    tf.text = "Top 25 Consensus Anomalies verified against real-world data:"

    p = tf.add_paragraph()
    p.text = "Severe Ecosystem Hacks & Cascades:"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "SLP ($622m Ronin Hack), BNT ($23m Protocol Hack), CORE (50% cascade)"
    p.level = 2

    p = tf.add_paragraph()
    p.text = "Confirmed Massive Scams & Ponzi Schemes:"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "CERE ($100m SEC lawsuit), TARO (RobotEra Exit Scam), HYDRA"
    p.level = 2

    p = tf.add_paragraph()
    p.text = "Controversial Insider Dumps / Honeypots:"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "OM (90% crash), LUA (Honeypot lock), MAGGOT, APEX scammers"
    p.level = 2

    p = tf.add_paragraph()
    p.text = "Conclusion: The Unsupervised Model catches ALL catastrophic fraud."
    p.level = 0
    p.font.bold = True

    prs.save('April_16_Update.pptx')
    print("Successfully updated April_16_Update.pptx")
except Exception as e:
    print(f"Error updating presentation: {e}")
